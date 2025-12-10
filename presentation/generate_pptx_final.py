#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génération automatique du PowerPoint final.
Version ULTRA PRO – UML + Benchmarks + Code + Design Automatique.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import cairosvg

# Pour générer les captures de code
try:
    from generate_code_snapshots import generate_all_snapshots
except ImportError:
    generate_all_snapshots = None

# ---------------------------------------------------------------------------
# Chemins de base
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent          # .../server_project/presentation
PPT_DIR = ROOT / "presentation"                 # .../presentation/presentation
BG_DIR = PPT_DIR / "backgrounds"

PROJECT_ROOT = ROOT.parents[0]                  # .../server_project
UML_DIR = PROJECT_ROOT / "docs" / "uml"
FIG_DIR = PROJECT_ROOT / "python" / "figures"
SNAP_DIR = ROOT / "code_snapshots"

OUTPUT = PPT_DIR / "presentation_finale_serveur.pptx"

TITLE_COLOR = RGBColor(20, 20, 20)
TEXT_COLOR = RGBColor(40, 40, 40)


# ---------------------------------------------------------------------------
# Helpers images
# ---------------------------------------------------------------------------

def svg_to_png(svg_path: Path) -> Path:
    """
    Convertit un fichier SVG en PNG (même nom, extension .png).
    Si le PNG existe déjà, on le réutilise.
    """
    png_path = svg_path.with_suffix(".png")
    if png_path.exists():
        return png_path

    if not svg_path.exists():
        print(f"[WARN] SVG introuvable : {svg_path}")
        return svg_path

    print(f"[SVG→PNG] {svg_path} -> {png_path}")
    cairosvg.svg2png(url=str(svg_path), write_to=str(png_path))
    return png_path


def resolve_image(path: Path) -> Path:
    """
    Résout automatiquement l'image à utiliser :
      - si c'est un SVG → converti en PNG
      - si c'est un PNG existant → utilisé tel quel
      - sinon, tente l'extension .png
    """
    if path.suffix.lower() == ".svg":
        return svg_to_png(path)

    if path.exists():
        return path

    alt = path.with_suffix(".png")
    if alt.exists():
        return alt

    print(f"[WARN] Image introuvable : {path}")
    return path


# ---------------------------------------------------------------------------
# Helpers slides
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, background: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(background),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

    tx = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(10), Inches(2)
    ).text_frame
    tx.text = title
    p = tx.paragraphs[0]
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    sub = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(10), Inches(1)
    ).text_frame
    sub.text = subtitle
    p2 = sub.paragraphs[0]
    p2.font.size = Pt(26)
    p2.font.color.rgb = TEXT_COLOR


def add_simple_image_slide(prs, title, image_path: Path, background: Path):
    """
    Slide simple : titre + image (UML / graph).
    """
    img = resolve_image(image_path)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(background),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(10), Inches(1))
    tx = t.text_frame
    tx.text = title
    p = tx.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    if img.exists():
        slide.shapes.add_picture(
            str(img),
            Inches(1),
            Inches(2),
            width=Inches(10),
        )
    else:
        # Fallback : texte d’avertissement
        warn_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(10), Inches(1.5)
        ).text_frame
        warn_box.text = f"[Image manquante] {img}"
        warn_box.paragraphs[0].font.size = Pt(18)
        warn_box.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)


def add_code_explained_slide(
    prs,
    title: str,
    description_lines,
    image_path: Path,
    background: Path,
):
    """
    Slide combinant :
      - un titre
      - une zone de texte 'fonctionnement'
      - une capture de code (PNG)
    """
    img = resolve_image(image_path)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(background),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

    # Titre
    t_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.5), Inches(10), Inches(0.8)
    )
    tf_title = t_title.text_frame
    tf_title.text = title
    p_title = tf_title.paragraphs[0]
    p_title.font.size = Pt(34)
    p_title.font.bold = True
    p_title.font.color.rgb = TITLE_COLOR

    # Description du fonctionnement
    t_desc = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.4), Inches(10), Inches(2)
    )
    tf_desc = t_desc.text_frame
    tf_desc.word_wrap = True

    first = True
    for line in description_lines:
        if first:
            tf_desc.text = line
            p = tf_desc.paragraphs[0]
            first = False
        else:
            p = tf_desc.add_paragraph()
            p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

    # Image de code
    if img.exists():
        slide.shapes.add_picture(
            str(img),
            Inches(0.7),
            Inches(3.0),
            width=Inches(10.5),
        )
    else:
        warn_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.0), Inches(10), Inches(1.5)
        ).text_frame
        warn_box.text = f"[Image code manquante] {img}"
        warn_box.paragraphs[0].font.size = Pt(18)
        warn_box.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)


# ---------------------------------------------------------------------------
# Génération PPT
# ---------------------------------------------------------------------------

def generate_ppt():
    # 1) Optionnel : génération des snapshots de code
    if generate_all_snapshots is not None:
        try:
            generate_all_snapshots()
        except Exception as e:
            print(f"[WARN] Génération snapshots code échouée : {e}")
    else:
        print("[WARN] Module generate_code_snapshots introuvable, pas de PNG code auto.")

    prs = Presentation()
    bg_light = BG_DIR / "bg_light.png"

    # Slide 1 : Titre
    add_title_slide(
        prs,
        "Serveurs TCP & HTTP Haute Performance",
        "Multi-thread | Queue FIFO | Benchmarks | C/POSIX",
        bg_light,
    )

    # 2) UML / Architecture / Threads
    uml_sections = [
        ("Architecture Globale", UML_DIR / "uml_architecture.svg"),
        ("Queue FIFO Thread-Safe", UML_DIR / "uml_queue.svg"),
        ("Threads & Workers", UML_DIR / "uml_threads.svg"),
        ("Séquence TCP Mono-thread", UML_DIR / "uml_seq_tcp_monothread.svg"),
        ("Séquence TCP Multi-thread", UML_DIR / "uml_seq_tcp_multithread.svg"),
        ("Séquence HTTP Mono-thread", UML_DIR / "uml_seq_http_monothread.svg"),
        ("Séquence HTTP Multi-thread", UML_DIR / "uml_seq_http_multithread.svg"),
    ]

    for title, img in uml_sections:
        add_simple_image_slide(prs, title, img, bg_light)

    # 3) Graphes de benchmark
    bench_sections = [
        ("Throughput (req/s)", FIG_DIR / "1-throughput.png"),
        ("Latence P99 (µs)", FIG_DIR / "2-latency_p99.png"),
        ("Utilisation CPU", FIG_DIR / "3-cpu.png"),
        ("Mémoire", FIG_DIR / "4-memory.png"),
        ("Speedup Multi-thread", FIG_DIR / "5-speedup.png"),
    ]

    for title, img in bench_sections:
        if img.exists():
            add_simple_image_slide(prs, title, img, bg_light)
        else:
            print(f"[WARN] Figure de benchmark manquante : {img}")

    # 4) Slides de CODE + EXPLICATIONS

    code_specs = [
        (
            "HTTP – Parser & Réponses (http.c)",
            SNAP_DIR / "code_http_c.png",
            [
                "Implémente le parsing de la ligne de requête HTTP (méthode, chemin, query).",
                "Gère un découpage robuste des espaces et des paramètres après '?'.",
                "Fournit une API simple pour les serveurs : parse_http_request() + send_http_response().",
                "Encapsule la construction d’une réponse HTTP 1.1 (status line, headers, body).",
            ],
        ),
        (
            "HTTP – Interface & Constantes (http.h)",
            SNAP_DIR / "code_http_h.png",
            [
                "Expose les prototypes du parser et de l’émetteur de réponse HTTP.",
                "Centralise les tailles de buffers et types utilisés côté HTTP.",
                "Permet de partager le même moteur HTTP entre serveur mono et multi-thread.",
            ],
        ),
        (
            "Queue FIFO Thread-Safe (queue.c)",
            SNAP_DIR / "code_queue_c.png",
            [
                "Implémente une file FIFO bornée, thread-safe, utilisée par le serveur multi-thread.",
                "Utilise un mutex + 2 variables de condition (not_empty / not_full).",
                "Supporte un mode shutdown propre pour réveiller tous les workers et le dispatcher.",
                "Assure un comportement strictement FIFO et évite les conditions de course.",
            ],
        ),
        (
            "Queue FIFO – Interface (queue.h)",
            SNAP_DIR / "code_queue_h.png",
            [
                "Définit la structure queue_t (head, tail, size, size_max, mutex, cond).",
                "Expose queue_init(), queue_push(), queue_pop(), queue_shutdown(), queue_destroy().",
                "Permet de réutiliser la même abstration pour TCP et HTTP (multi-thread).",
            ],
        ),
        (
            "Serveur TCP Mono-thread (serveur_mono.c)",
            SNAP_DIR / "code_serveur_mono_c.png",
            [
                "Boucle accept() → recv() → traitement_lourd() → send() pour un seul client à la fois.",
                "Utilise un traitement CPU-bound simulé (~100ms) pour mesurer la saturation.",
                "Renvoie le carré du nombre reçu (+ timestamp µs) au client.",
                "SIGINT handler simple : fermeture du socket serveur et exit immédiat.",
            ],
        ),
        (
            "Serveur HTTP Mono-thread (serveur_mono_http.c)",
            SNAP_DIR / "code_serveur_mono_http_c.png",
            [
                "Accepte les connexions une par une sur le port HTTP mono-thread (8080).",
                "Parse la requête brute via http.c, route vers /, /hello, /time, /stats.",
                "Gère des timeouts recv() pour éviter les connexions bloquées.",
                "Idéal comme référence séquentielle pour comparer au multi-thread HTTP.",
            ],
        ),
        (
            "Serveur TCP Multi-thread (serveur_multi.c)",
            SNAP_DIR / "code_serveur_multi_c.png",
            [
                "Crée un pool fixe de WORKER_COUNT threads dès le démarrage.",
                "Le thread principal accepte les connexions et les pousse dans la queue FIFO.",
                "Chaque worker dépile un fd, exécute traitement_lourd(), renvoie la réponse, ferme le fd.",
                "Gère SIGINT + queue_shutdown() pour un arrêt propre sans deadlock.",
            ],
        ),
        (
            "Serveur HTTP Multi-thread (serveur_multi_http.c)",
            SNAP_DIR / "code_serveur_multi_http_c.png",
            [
                "Architecture identique à TCP multi-thread, mais au niveau HTTP 1.1 (port 8081).",
                "Workers parse la requête HTTP, appellent route_request(), renvoient une réponse JSON/HTML.",
                "Statistiques globales /stats protégées par mutex (total_requests, hello_requests, 404).",
                "Utilise SO_RCVTIMEO pour limiter la durée de blocage sur recv().",
            ],
        ),
        (
            "Client de Charge / Benchmarks (python/client_stress.py)",
            SNAP_DIR / "code_client_stress_py.png",
            [
                "Génère des centaines de clients concurrents pour mesurer throughput et latence.",
                "Ouvre des connexions TCP/HTTP, envoie des requêtes, collecte les temps de réponse.",
                "Produit des métriques agrégées (P50, P95, P99, RPS) en JSON / Excel.",
                "Alimente le dashboard Plotly + les figures utilisées dans la présentation.",
            ],
        ),
    ]

    for title, img, desc in code_specs:
        add_code_explained_slide(prs, title, desc, img, bg_light)

    # Sauvegarde
    PPT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"✔ PowerPoint généré : {OUTPUT}")


if __name__ == "__main__":
    generate_ppt()

